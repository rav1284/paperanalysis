from pptx import Presentation

def create_ppt(summary):
    prs = Presentation()

from pptx import Presentation

def create_ppt(summary):
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = summary["title"]
    slide.placeholders[1].text = "Research Paper Summary"

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Abstract"
    slide.placeholders[1].text = summary["abstract"]

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = summary["introduction"]

    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion"
    slide.placeholders[1].text = summary["conclusion"]

    prs.save("output.pptx")